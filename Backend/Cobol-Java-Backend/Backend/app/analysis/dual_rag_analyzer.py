import os
import json
import logging
import hashlib
import re
from datetime import datetime
from typing import Dict, List, Any, Optional
from dataclasses import dataclass
from enum import Enum
from pathlib import Path

# Document processing
import docx
import PyPDF2
from pptx import Presentation

# RAG components
from langchain_community.vectorstores import FAISS
from langchain_openai import AzureOpenAIEmbeddings
from langchain_openai import AzureChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_text_splitters import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)

class COBOLNodeType(Enum):
    PROGRAM = "program"
    DIVISION = "division"
    SECTION = "section"
    DATA_ITEM = "data_item"
    STATEMENT = "statement"
    FILE_DESCRIPTOR = "file_descriptor"

@dataclass
class COBOLNode:
    node_type: COBOLNodeType
    name: str
    content: str
    line_number: int
    children: List['COBOLNode']
    attributes: Dict[str, Any]
    dependencies: List[str] = None

class DualRAGCOBOLAnalyzer:
    """
    Enhanced COBOL analyzer with dual RAG system:
    1. Standards RAG - for client documents/standards
    2. Project RAG - for project patterns and connections
    """
    
    def __init__(self, azure_config: Dict[str, str], output_dir: str = "cobol_analysis"):
        self.output_dir = output_dir
        self.documents_dir = "documents"  # Standards documents folder
        
        # Output files
        self.analysis_file = os.path.join(output_dir, "project_analysis.json")
        self.connections_file = os.path.join(output_dir, "file_connections.json")
        
        # RAG directories - SEPARATE STORES
        self.standards_rag_dir = os.path.join(output_dir, "standards_rag")
        self.project_rag_dir = os.path.join(output_dir, "project_rag")
        
        # Create all directories
        for directory in [output_dir, self.documents_dir, self.standards_rag_dir, self.project_rag_dir]:
            os.makedirs(directory, exist_ok=True)
        
        # Initialize Azure components
        self.embeddings = AzureOpenAIEmbeddings(
            azure_endpoint=azure_config["AZURE_OPENAI_EMBED_API_ENDPOINT"],
            api_key=azure_config["AZURE_OPENAI_EMBED_API_KEY"],
            api_version=azure_config["AZURE_OPENAI_EMBED_VERSION"],
            model=azure_config["AZURE_OPENAI_EMBED_MODEL"]
        )
        
        self.llm = AzureChatOpenAI(
            azure_endpoint=azure_config["AZURE_OPENAI_ENDPOINT"],
            api_key=azure_config["AZURE_OPENAI_API_KEY"],
            api_version=azure_config["AZURE_OPENAI_API_VERSION"],
            deployment_name=azure_config["AZURE_OPENAI_DEPLOYMENT"],
            temperature=0.2
        )
        
        # Text splitter for documents
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        
        # Dual RAG stores
        self.standards_vector_store = None    # For standards documents
        self.project_vector_store = None      # For project patterns
        
        # Project analysis storage
        self.project_analysis = {
            "files": {},
            "connections": {},
            "dependencies": {},
            "field_mappings": {},
            "business_entities": [],
            "conversion_patterns": []
        }
        
        # Load existing RAG stores
        self._load_existing_rag_stores()
        
        logger.info("🚀 Dual RAG COBOL Analyzer initialized")
        logger.info(f"📚 Standards RAG: {self.standards_rag_dir}")
        logger.info(f"🔗 Project RAG: {self.project_rag_dir}")
        logger.info(f"📄 Documents folder: {self.documents_dir}")
    
    def _load_existing_rag_stores(self):
        """Load existing RAG stores if available"""
        
        # Load Standards RAG
        standards_index_path = os.path.join(self.standards_rag_dir, "index.faiss")
        if os.path.exists(standards_index_path):
            try:
                self.standards_vector_store = FAISS.load_local(
                    self.standards_rag_dir,
                    self.embeddings,
                    allow_dangerous_deserialization=True
                )
                logger.info("✅ Loaded existing Standards RAG")
            except Exception as e:
                logger.warning(f"⚠️ Could not load Standards RAG: {e}")
        
        # Load Project RAG
        project_index_path = os.path.join(self.project_rag_dir, "index.faiss")
        if os.path.exists(project_index_path):
            try:
                self.project_vector_store = FAISS.load_local(
                    self.project_rag_dir,
                    self.embeddings,
                    allow_dangerous_deserialization=True
                )
                logger.info("✅ Loaded existing Project RAG")
            except Exception as e:
                logger.warning(f"⚠️ Could not load Project RAG: {e}")
    
    def process_standards_documents(self):
        """Process all documents in documents/ folder and store in Standards RAG"""
        
        logger.info(f"📄 Processing standards documents from {self.documents_dir}")
        
        if not os.path.exists(self.documents_dir):
            logger.warning(f"⚠️ Documents directory {self.documents_dir} not found")
            return
        
        all_documents = []
        documents_metadata = []
        
        # Process each document in the documents folder
        for filename in os.listdir(self.documents_dir):
            file_path = os.path.join(self.documents_dir, filename)
            
            try:
                logger.info(f"📖 Processing {filename}")
                
                # Extract text based on file type
                extracted_text = self._extract_text_from_document(file_path)
                
                if extracted_text:
                    # Split text into chunks
                    chunks = self.text_splitter.split_text(extracted_text)
                    
                    # Add each chunk with metadata
                    for i, chunk in enumerate(chunks):
                        all_documents.append(chunk)
                        documents_metadata.append({
                            "source_file": filename,
                            "chunk_id": i,
                            "document_type": "standards",
                            "file_type": Path(filename).suffix.lower(),
                            "processed_at": datetime.now().isoformat()
                        })
                    
                    logger.info(f"✅ Processed {filename}: {len(chunks)} chunks")
                else:
                    logger.warning(f"⚠️ No text extracted from {filename}")
                    
            except Exception as e:
                logger.error(f"❌ Error processing {filename}: {e}")
        
        # Build/Update Standards RAG
        if all_documents:
            if self.standards_vector_store is None:
                # Create new vector store
                self.standards_vector_store = FAISS.from_texts(
                    all_documents, 
                    self.embeddings, 
                    metadatas=documents_metadata
                )
                logger.info(f"🆕 Created new Standards RAG with {len(all_documents)} chunks")
            else:
                # Add to existing vector store
                self.standards_vector_store.add_texts(all_documents, metadatas=documents_metadata)
                logger.info(f"➕ Added {len(all_documents)} chunks to existing Standards RAG")
            
            # Save Standards RAG
            self.standards_vector_store.save_local(self.standards_rag_dir)
            
            # Save metadata
            standards_info = {
                "total_chunks": len(all_documents),
                "processed_files": list(set(meta["source_file"] for meta in documents_metadata)),
                "last_updated": datetime.now().isoformat()
            }
            
            with open(os.path.join(self.standards_rag_dir, "documents.json"), 'w') as f:
                json.dump(standards_info, f, indent=2)
            
            logger.info(f"💾 Standards RAG saved to {self.standards_rag_dir}")
        else:
            logger.warning("⚠️ No documents processed for Standards RAG")
    
    def _extract_text_from_document(self, file_path: str) -> str:
        """Extract text from various document formats"""
        
        file_ext = Path(file_path).suffix.lower()
        
        try:
            if file_ext == '.txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            
            elif file_ext == '.pdf':
                text = ""
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                return text
            
            elif file_ext == '.docx':
                doc = docx.Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text
            
            elif file_ext == '.pptx':
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            
            else:
                logger.warning(f"⚠️ Unsupported file type: {file_ext}")
                return ""
                
        except Exception as e:
            logger.error(f"❌ Error extracting text from {file_path}: {e}")
            return ""
    
    def analyze_project(self, project_files: Dict[str, str]) -> Dict[str, Any]:
        """
        Main analysis function - analyzes entire project and builds Project RAG
        """
        logger.info(f"🔍 Analyzing project with {len(project_files)} files")
        
        # Step 1: Process standards documents first
        self.process_standards_documents()
        
        # Step 2: Analyze project files
        for filename, content in project_files.items():
            self._analyze_single_file(filename, content)
        
        # Step 3: Build connections between files
        self._build_file_connections()
        
        # Step 4: Extract business entities and field mappings
        self._extract_business_entities()
        
        # Step 5: Build conversion patterns
        self._build_conversion_patterns()
        
        # Step 6: Store project data in Project RAG
        self._build_project_rag()
        
        # Step 7: Save analysis
        self._save_analysis()
        
        logger.info("✅ Project analysis completed")
        return self.project_analysis
    
    def _build_project_rag(self):
        """Build Project RAG from analysis data"""
        
        logger.info("🔗 Building Project RAG from analysis data")
        
        project_texts = []
        project_metadata = []
        
        # 1. File connections as text
        for filename, conn_info in self.project_analysis["connections"].items():
            if conn_info["depends_on"]:
                connection_text = f"File {filename} depends on: {', '.join(conn_info['depends_on'])}"
                project_texts.append(connection_text)
                project_metadata.append({
                    "type": "file_connection",
                    "source_file": filename,
                    "connection_count": len(conn_info["depends_on"]),
                    "document_type": "project_data"
                })
        
        # 2. Business entities as text
        for entity in self.project_analysis["business_entities"]:
            entity_text = f"Business entity {entity['name']} from {entity['source_file']} maps to Java class {entity['java_class_name']}"
            project_texts.append(entity_text)
            project_metadata.append({
                "type": "business_entity",
                "entity_name": entity["name"],
                "source_file": entity["source_file"],
                "java_class": entity["java_class_name"],
                "document_type": "project_data"
            })
        
        # 3. Conversion patterns as text
        for pattern in self.project_analysis["conversion_patterns"]:
            pattern_text = f"{pattern['pattern_type']}: {pattern['cobol_pattern']} converts to {pattern['java_equivalent']}"
            project_texts.append(pattern_text)
            project_metadata.append({
                "type": "conversion_pattern",
                "pattern_type": pattern["pattern_type"],
                "source_file": pattern["source_file"],
                "document_type": "project_data"
            })
        
        # 4. Field mappings as text
        for cobol_field, mapping in self.project_analysis["field_mappings"].items():
            mapping_text = f"COBOL field {cobol_field} with PIC {mapping['pic_clause']} maps to Java {mapping['java_type']} field {mapping['java_name']}"
            project_texts.append(mapping_text)
            project_metadata.append({
                "type": "field_mapping",
                "cobol_field": cobol_field,
                "java_field": mapping["java_name"],
                "java_type": mapping["java_type"],
                "document_type": "project_data"
            })
        
        # Build/Update Project RAG
        if project_texts:
            if self.project_vector_store is None:
                # Create new Project RAG
                self.project_vector_store = FAISS.from_texts(
                    project_texts,
                    self.embeddings,
                    metadatas=project_metadata
                )
                logger.info(f"🆕 Created new Project RAG with {len(project_texts)} entries")
            else:
                # Add to existing Project RAG
                self.project_vector_store.add_texts(project_texts, metadatas=project_metadata)
                logger.info(f"➕ Added {len(project_texts)} entries to existing Project RAG")
            
            # Save Project RAG
            self.project_vector_store.save_local(self.project_rag_dir)
            
            # Save metadata
            project_rag_info = {
                "total_entries": len(project_texts),
                "entry_types": {
                    "file_connections": len([m for m in project_metadata if m["type"] == "file_connection"]),
                    "business_entities": len([m for m in project_metadata if m["type"] == "business_entity"]),
                    "conversion_patterns": len([m for m in project_metadata if m["type"] == "conversion_pattern"]),
                    "field_mappings": len([m for m in project_metadata if m["type"] == "field_mapping"])
                },
                "last_updated": datetime.now().isoformat()
            }
            
            with open(os.path.join(self.project_rag_dir, "patterns.json"), 'w') as f:
                json.dump(project_rag_info, f, indent=2)
            
            logger.info(f"💾 Project RAG saved to {self.project_rag_dir}")
    
    def get_standards_context(self, query: str, k: int = 3) -> str:
        """Get context from Standards RAG"""
        
        if not self.standards_vector_store:
            return "No standards documents available."
        
        try:
            similar_docs = self.standards_vector_store.similarity_search(query, k=k)
            
            context_parts = ["Standards and Guidelines Context:\n"]
            for i, doc in enumerate(similar_docs, 1):
                source_file = doc.metadata.get("source_file", "unknown")
                context_parts.append(f"{i}. From {source_file}:")
                context_parts.append(f"   {doc.page_content[:300]}...")
                context_parts.append("")
            
            return "\n".join(context_parts)
            
        except Exception as e:
            logger.warning(f"⚠️ Error getting standards context: {e}")
            return "Error retrieving standards context."
    
    def get_project_context(self, query: str, k: int = 5) -> str:
        """Get context from Project RAG"""
        
        if not self.project_vector_store:
            return "No project patterns available."
        
        try:
            similar_docs = self.project_vector_store.similarity_search(query, k=k)
            
            context_parts = ["Project Patterns and Connections Context:\n"]
            for i, doc in enumerate(similar_docs, 1):
                doc_type = doc.metadata.get("type", "unknown")
                context_parts.append(f"{i}. {doc_type.title()}: {doc.page_content}")
                context_parts.append("")
            
            return "\n".join(context_parts)
            
        except Exception as e:
            logger.warning(f"⚠️ Error getting project context: {e}")
            return "Error retrieving project context."
    
    def get_combined_conversion_context(self, cobol_code: str) -> str:
        """Get combined context from both RAG stores"""
        
        # Get context from both RAG stores
        standards_context = self.get_standards_context(cobol_code, k=2)
        project_context = self.get_project_context(cobol_code, k=3)
        
        # Combine contexts
        combined_context = f"""
CONVERSION CONTEXT

{standards_context}

{project_context}

Project Overview:
- Total Files: {len(self.project_analysis['files'])}
- Business Entities: {len(self.project_analysis['business_entities'])}
- File Connections: {sum(len(conn['depends_on']) for conn in self.project_analysis['connections'].values())}
- Conversion Patterns: {len(self.project_analysis['conversion_patterns'])}
"""
        
        return combined_context
    
    # [Keep all the existing methods from the original SimplifiedCOBOLAnalyzer]
    # _analyze_single_file, _extract_dependencies, _build_file_connections, 
    # _extract_business_entities, _build_conversion_patterns, etc.
    
    def _analyze_single_file(self, filename: str, content: str):
        """Analyze single file for structure and dependencies"""
        
        file_analysis = {
            "filename": filename,
            "file_type": self._determine_file_type(filename),
            "size": len(content),
            "lines": len(content.split('\n')),
            "program_id": None,
            "divisions": [],
            "data_items": [],
            "dependencies": {
                "copy_books": [],
                "called_programs": [],
                "file_references": [],
                "vsam_files": [],
                "jcl_references": []
            },
            "business_logic": [],
            "sql_statements": []
        }
        
        lines = content.split('\n')
        
        # Extract program structure
        for i, line in enumerate(lines, 1):
            line_clean = line.strip()
            if not line_clean or line_clean.startswith('*'):
                continue
            
            # Program ID
            if not file_analysis["program_id"]:
                prog_match = re.search(r'PROGRAM-ID\.\s*([A-Z0-9-]+)', line_clean, re.I)
                if prog_match:
                    file_analysis["program_id"] = prog_match.group(1)
            
            # Divisions
            div_match = re.search(r'(IDENTIFICATION|ENVIRONMENT|DATA|PROCEDURE)\s+DIVISION', line_clean, re.I)
            if div_match:
                file_analysis["divisions"].append({
                    "name": div_match.group(1),
                    "line": i
                })
            
            # Data items (01 level)
            data_match = re.search(r'^\s*01\s+([A-Z0-9-]+)', line_clean, re.I)
            if data_match:
                pic_match = re.search(r'PIC\s+([9AXV\(\)\.]+)', line_clean, re.I)
                file_analysis["data_items"].append({
                    "name": data_match.group(1),
                    "line": i,
                    "pic": pic_match.group(1) if pic_match else None,
                    "java_type": self._determine_java_type(pic_match.group(1) if pic_match else None)
                })
            
            # Dependencies
            self._extract_dependencies(line_clean, file_analysis["dependencies"])
            
            # Business logic patterns
            if any(keyword in line_clean.upper() for keyword in ['PERFORM', 'IF', 'COMPUTE', 'MOVE', 'CALL']):
                file_analysis["business_logic"].append({
                    "line": i,
                    "statement": line_clean[:100],  # First 100 chars
                    "type": self._classify_statement(line_clean)
                })
            
            # SQL statements
            if 'EXEC SQL' in line_clean.upper():
                file_analysis["sql_statements"].append({
                    "line": i,
                    "statement": line_clean
                })
        
        self.project_analysis["files"][filename] = file_analysis
    
    def _extract_dependencies(self, line: str, deps: Dict[str, List]):
        """Extract all types of dependencies from a line"""
        line_upper = line.upper()
        
        # Copy books
        copy_match = re.search(r'COPY\s+([A-Z0-9-]+)', line_upper)
        if copy_match:
            deps["copy_books"].append(copy_match.group(1))
        
        # Called programs
        call_match = re.search(r'CALL\s+[\'"]([^\'\"]+)[\'"]', line_upper)
        if call_match:
            deps["called_programs"].append(call_match.group(1))
        
        # File references
        if any(keyword in line_upper for keyword in ['SELECT', 'ASSIGN', 'FD']):
            file_match = re.search(r'(SELECT|FD)\s+([A-Z0-9-]+)', line_upper)
            if file_match:
                deps["file_references"].append(file_match.group(2))
        
        # VSAM file references
        if 'VSAM' in line_upper or 'KSDS' in line_upper or 'ESDS' in line_upper:
            vsam_match = re.search(r'([A-Z0-9-]+).*VSAM', line_upper)
            if vsam_match:
                deps["vsam_files"].append(vsam_match.group(1))
    
    def _build_file_connections(self):
        """Build connections between files based on dependencies"""
        
        connections = {}
        
        for filename, file_info in self.project_analysis["files"].items():
            file_connections = {
                "depends_on": [],
                "used_by": [],
                "connection_type": []
            }
            
            deps = file_info["dependencies"]
            
            # Check connections to other files in project
            for other_filename in self.project_analysis["files"].keys():
                if other_filename == filename:
                    continue
                
                other_base = os.path.splitext(other_filename)[0]
                
                # Check if this file depends on the other
                if (other_base in deps["copy_books"] or 
                    other_base in deps["called_programs"] or
                    other_base in deps["file_references"]):
                    
                    file_connections["depends_on"].append(other_filename)
                    
                    # Determine connection type
                    if other_base in deps["copy_books"]:
                        file_connections["connection_type"].append(f"{other_filename}:COPY")
                    if other_base in deps["called_programs"]:
                        file_connections["connection_type"].append(f"{other_filename}:CALL")
                    if other_base in deps["file_references"]:
                        file_connections["connection_type"].append(f"{other_filename}:FILE")
            
            connections[filename] = file_connections
        
        # Build reverse connections (used_by)
        for filename, conn_info in connections.items():
            for dep_file in conn_info["depends_on"]:
                if dep_file in connections:
                    connections[dep_file]["used_by"].append(filename)
        
        self.project_analysis["connections"] = connections
    
    def _extract_business_entities(self):
        """Extract business entities from data structures"""
        
        entities = []
        field_mappings = {}
        
        for filename, file_info in self.project_analysis["files"].items():
            for data_item in file_info["data_items"]:
                # Consider 01-level items as potential entities
                entity = {
                    "name": data_item["name"],
                    "source_file": filename,
                    "line": data_item["line"],
                    "java_class_name": self._to_java_class_name(data_item["name"]),
                    "fields": []
                }
                
                # Look for child fields (this is simplified)
                entity["estimated_complexity"] = "medium" if len(file_info["data_items"]) > 10 else "simple"
                
                entities.append(entity)
                
                # Field mapping
                field_mappings[data_item["name"]] = {
                    "cobol_name": data_item["name"],
                    "java_name": self._to_java_field_name(data_item["name"]),
                    "java_type": data_item["java_type"],
                    "pic_clause": data_item["pic"]
                }
        
        self.project_analysis["business_entities"] = entities
        self.project_analysis["field_mappings"] = field_mappings
    
    def _build_conversion_patterns(self):
        """Build simple conversion patterns for RAG"""
        
        patterns = []
        
        # Pattern 1: Data structure patterns
        for entity in self.project_analysis["business_entities"]:
            pattern = {
                "pattern_type": "data_structure",
                "cobol_pattern": f"01 {entity['name']}",
                "java_equivalent": f"public class {entity['java_class_name']}",
                "context": f"Convert COBOL record {entity['name']} to Java class",
                "source_file": entity["source_file"]
            }
            patterns.append(pattern)
        
        # Pattern 2: File operation patterns
        for filename, file_info in self.project_analysis["files"].items():
            if file_info["dependencies"]["file_references"]:
                for file_ref in file_info["dependencies"]["file_references"]:
                    pattern = {
                        "pattern_type": "file_operation",
                        "cobol_pattern": f"FD {file_ref}",
                        "java_equivalent": f"// Repository pattern for {file_ref}",
                        "context": f"Convert file operations to repository pattern",
                        "source_file": filename
                    }
                    patterns.append(pattern)
        
        # Pattern 3: Business logic patterns
        for filename, file_info in self.project_analysis["files"].items():
            for logic in file_info["business_logic"]:
                if logic["type"] in ["PERFORM", "IF", "COMPUTE"]:
                    pattern = {
                        "pattern_type": "business_logic",
                        "cobol_pattern": logic["statement"],
                        "java_equivalent": f"// {logic['type']} converted to Java",
                        "context": f"Convert {logic['type']} statement to Java",
                        "source_file": filename
                    }
                    patterns.append(pattern)
        
        self.project_analysis["conversion_patterns"] = patterns
    
    def smart_chunk_if_needed(self, content: str, max_size: int = 30000) -> List[str]:
        """Simple chunking for large files"""
        
        if len(content) <= max_size:
            return [content]
        
        chunks = []
        lines = content.split('\n')
        current_chunk = []
        current_size = 0
        
        for line in lines:
            if current_size + len(line) > max_size and current_chunk:
                chunks.append('\n'.join(current_chunk))
                current_chunk = [line]
                current_size = len(line)
            else:
                current_chunk.append(line)
                current_size += len(line)
        
        if current_chunk:
            chunks.append('\n'.join(current_chunk))
        
        return chunks
    
    def _save_analysis(self):
        """Save analysis to files"""
        
        # Main analysis file
        with open(self.analysis_file, 'w', encoding='utf-8') as f:
            json.dump(self.project_analysis, f, indent=2, ensure_ascii=False)
        
        # Simplified connections file
        connections_summary = {
            "project_summary": {
                "total_files": len(self.project_analysis["files"]),
                "total_connections": sum(len(conn["depends_on"]) for conn in self.project_analysis["connections"].values()),
                "business_entities": len(self.project_analysis["business_entities"]),
                "conversion_patterns": len(self.project_analysis["conversion_patterns"])
            },
            "file_connections": self.project_analysis["connections"],
            "dependency_matrix": self._build_dependency_matrix(),
            "rag_summary": {
                "standards_rag_available": self.standards_vector_store is not None,
                "project_rag_available": self.project_vector_store is not None,
                "standards_rag_path": self.standards_rag_dir,
                "project_rag_path": self.project_rag_dir
            }
        }
        
        with open(self.connections_file, 'w', encoding='utf-8') as f:
            json.dump(connections_summary, f, indent=2, ensure_ascii=False)
        
        # Human-readable summary
        self._save_readable_summary()
        
        logger.info(f"💾 Analysis saved to {self.output_dir}")
    
    def _build_dependency_matrix(self) -> Dict[str, List[str]]:
        """Build simple dependency matrix"""
        matrix = {}
        for filename, conn_info in self.project_analysis["connections"].items():
            matrix[filename] = conn_info["depends_on"]
        return matrix
    
    def _save_readable_summary(self):
        """Save human-readable summary"""
        
        summary_file = os.path.join(self.output_dir, "project_summary.md")
        
        with open(summary_file, 'w', encoding='utf-8') as f:
            f.write("# COBOL Project Analysis Summary with Dual RAG\n\n")
            f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
            
            # Project overview
            f.write("## Project Overview\n\n")
            f.write(f"- **Total Files:** {len(self.project_analysis['files'])}\n")
            f.write(f"- **Business Entities:** {len(self.project_analysis['business_entities'])}\n")
            f.write(f"- **Total Connections:** {sum(len(conn['depends_on']) for conn in self.project_analysis['connections'].values())}\n\n")
            
            # RAG Status
            f.write("## RAG System Status\n\n")
            f.write(f"- **Standards RAG:** {'✅ Active' if self.standards_vector_store else '❌ Not Available'}\n")
            f.write(f"- **Project RAG:** {'✅ Active' if self.project_vector_store else '❌ Not Available'}\n")
            f.write(f"- **Standards Location:** `{self.standards_rag_dir}`\n")
            f.write(f"- **Project Location:** `{self.project_rag_dir}`\n\n")
            
            # File connections
            f.write("## File Dependencies\n\n")
            for filename, conn_info in self.project_analysis["connections"].items():
                if conn_info["depends_on"]:
                    f.write(f"### {filename}\n")
                    f.write("**Depends on:**\n")
                    for dep in conn_info["depends_on"]:
                        f.write(f"- {dep}\n")
                    f.write("\n")
            
            # Business entities
            f.write("## Business Entities\n\n")
            for entity in self.project_analysis["business_entities"][:10]:  # Show first 10
                f.write(f"- **{entity['name']}** → `{entity['java_class_name']}` (from {entity['source_file']})\n")
            
            if len(self.project_analysis["business_entities"]) > 10:
                f.write(f"- ... and {len(self.project_analysis['business_entities']) - 10} more\n")
    
    # Helper methods (keep all existing ones)
    def _determine_file_type(self, filename: str) -> str:
        ext = os.path.splitext(filename)[1].lower()
        type_map = {
            '.cbl': 'cobol_program',
            '.cob': 'cobol_program', 
            '.cpy': 'copybook',
            '.inc': 'include',
            '.jcl': 'job_control',
            '.job': 'job_control',
            '.ctl': 'control_file'
        }
        return type_map.get(ext, 'unknown')
    
    def _classify_statement(self, line: str) -> str:
        line_upper = line.upper()
        if line_upper.startswith('PERFORM'):
            return 'PERFORM'
        elif line_upper.startswith('IF'):
            return 'IF'
        elif line_upper.startswith('MOVE'):
            return 'MOVE'
        elif line_upper.startswith('COMPUTE'):
            return 'COMPUTE'
        elif line_upper.startswith('CALL'):
            return 'CALL'
        else:
            return 'OTHER'
    
    def _determine_java_type(self, pic_clause: str) -> str:
        if not pic_clause:
            return "String"
        pic_upper = pic_clause.upper()
        if 'X' in pic_upper or 'A' in pic_upper:
            return "String"
        elif '9' in pic_upper:
            if 'V' in pic_upper or '.' in pic_upper:
                return "BigDecimal"
            else:
                digit_count = pic_upper.count('9')
                return "Integer" if digit_count <= 9 else "Long"
        return "String"
    
    def _to_java_class_name(self, cobol_name: str) -> str:
        parts = cobol_name.replace('-', '_').split('_')
        return ''.join(word.capitalize() for word in parts)
    
    def _to_java_field_name(self, cobol_name: str) -> str:
        parts = cobol_name.replace('-', '_').split('_')
        return parts[0].lower() + ''.join(word.capitalize() for word in parts[1:])


